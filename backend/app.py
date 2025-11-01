from flask import Flask, request, jsonify, send_from_directory
from flask_cors import CORS
from langchain_community.vectorstores import Chroma
from langchain_community.embeddings import HuggingFaceEmbeddings
from langchain.text_splitter import RecursiveCharacterTextSplitter
import pandas as pd
import matplotlib.pyplot as plt
import numpy as np
import google.generativeai as genai
import os, uuid, warnings, shutil, re, json
from werkzeug.utils import secure_filename
from pptx import Presentation
import base64
import logging
import seaborn as sns
import requests
from bs4 import BeautifulSoup

# Configure matplotlib for headless environment
plt.switch_backend('Agg')
plt.style.use('dark_background')

# Configure logging
logging.basicConfig(level=logging.INFO)
logger = logging.getLogger(__name__)

# Configuration
GEMINI_API_KEY = os.getenv("GEMINI_API_KEY")
if GEMINI_API_KEY:
    genai.configure(api_key=GEMINI_API_KEY)
    logger.info("✓ Gemini API configured")
else:
    warnings.warn("⚠️ GEMINI_API_KEY not set")

# Document Loaders
LOADERS = {}
try:
    from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
    LOADERS.update({'pdf': PyPDFLoader, 'txt': TextLoader, 'docx': Docx2txtLoader})
    logger.info("✓ Document loaders loaded successfully")
except ImportError as e:
    logger.warning("⚠️ Limited loader support")

# Flask setup
app = Flask(__name__)
CORS(app)

# Directories
UPLOAD_FOLDER = os.path.join(os.path.dirname(__file__), 'uploads')
PLOTS_FOLDER = os.path.join(os.path.dirname(__file__), 'static', 'plots')
VECTOR_DB_DIR = os.path.join(os.path.dirname(__file__), 'db_miniLM')

# Ensure directories exist
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PLOTS_FOLDER, exist_ok=True)
os.makedirs(os.path.dirname(PLOTS_FOLDER), exist_ok=True)  # Ensure static folder exists

# Constants
EMBEDDING_MODEL = "sentence-transformers/all-MiniLM-L6-v2"
ALLOWED_EXTENSIONS = {'txt', 'pdf', 'docx', 'pptx', 'ppt', 'xlsx', 'xls', 'csv'}
excel_data_store = {}

# Keywords
GRAPH_KEYWORDS = ['graph', 'chart', 'plot', 'visualize', 'visualization', 'bar', 'line', 'pie', 'scatter', 'histogram', 'show', 'draw', 'create chart']
DOCUMENT_KEYWORDS = ['summarise', 'summarize', 'summary', 'document', 'content', 'explain', 'analyze', 'findings']

# Helper Functions
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in ALLOWED_EXTENSIONS

def clear_vector_database():
    try:
        if os.path.exists(VECTOR_DB_DIR):
            shutil.rmtree(VECTOR_DB_DIR)
            logger.info("🗑️ Cleared vector database")
            return True
    except Exception as e:
        logger.error(f"⚠️ Could not clear vector database: {str(e)}")
        return False
    return True

def has_documents():
    try:
        return os.path.exists(VECTOR_DB_DIR) and len(Chroma(persist_directory=VECTOR_DB_DIR, embedding_function=HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)).similarity_search("test", k=1)) > 0
    except:
        return False

def has_data():
    return len(excel_data_store) > 0

def is_graph_request(question):
    return any(keyword in question.lower() for keyword in GRAPH_KEYWORDS)

def is_document_specific_question(question):
    return any(keyword in question.lower() for keyword in DOCUMENT_KEYWORDS)

def search_web(query, max_results=3):
    """Simple web search using DuckDuckGo (you can replace with your preferred search API)"""
    try:
        # This is a simple implementation - replace with your preferred search API
        search_url = f"https://duckduckgo.com/html/?q={query}"
        headers = {
            'User-Agent': 'Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36 (KHTML, like Gecko) Chrome/91.0.4472.124 Safari/537.36'
        }
        
        response = requests.get(search_url, headers=headers, timeout=10)
        soup = BeautifulSoup(response.content, 'html.parser')
        
        results = []
        result_links = soup.find_all('a', class_='result__a')[:max_results]
        
        for link in result_links:
            title = link.get_text()
            url = link.get('href')
            if url and title:
                results.append({'title': title, 'url': url})
        
        return results
    except Exception as e:
        logger.error(f"Web search error: {str(e)}")
        return []

def ask_gemini(prompt, temp=0.7, include_web_context=False, web_query=None):
    if not GEMINI_API_KEY:
        return "⚠️ Gemini API key not set"
    
    try:
        model = genai.GenerativeModel("gemini-1.5-flash")
        
        # If web search is requested, add web context
        if include_web_context and web_query:
            web_results = search_web(web_query)
            if web_results:
                web_context = "Here are some relevant web search results:\n"
                for i, result in enumerate(web_results, 1):
                    web_context += f"{i}. {result['title']} - {result['url']}\n"
                
                prompt = f"{web_context}\n\nBased on current information and the above sources, {prompt}"
        
        response = model.generate_content(prompt, generation_config=genai.types.GenerationConfig(max_output_tokens=2000, temperature=temp))
        return response.text.strip() if response.text else "No response generated"
    except Exception as e:
        logger.error(f"Gemini API error: {str(e)}")
        return f"Error: {str(e)}"

def extract_ppt_content(filepath):
    """Extract PowerPoint content using python-pptx and Gemini for analysis"""
    try:
        if not os.path.exists(filepath):
            return f"Error: File not found at {filepath}"
        
        file_size = os.path.getsize(filepath)
        if file_size > 50 * 1024 * 1024:  # 50MB limit
            return "Error: File too large (>50MB)"
        
        logger.info(f"📄 Processing PowerPoint file: {filepath} ({file_size} bytes)")
        
        try:
            prs = Presentation(filepath)
        except Exception as e:
            return f"Error: Unable to open PowerPoint file - {str(e)}"
        
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            slide_text.append(f"--- Slide {i+1} ---")
            
            text_found = False
            for shape in slide.shapes:
                try:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        text_found = True
                    
                    if hasattr(shape, "table"):
                        table_text = []
                        for row in shape.table.rows:
                            row_text = []
                            for cell in row.cells:
                                if cell.text.strip():
                                    row_text.append(cell.text.strip())
                            if row_text:
                                table_text.append(" | ".join(row_text))
                        if table_text:
                            slide_text.extend(table_text)
                            text_found = True
                            
                except Exception as e:
                    logger.warning(f"Error processing shape on slide {i+1}: {str(e)}")
                    continue
            
            if not text_found:
                slide_text.append("[No text content found in this slide]")
                
            slides_content.append("\n".join(slide_text))
        
        if not slides_content:
            return "Error: No slides found in PowerPoint file"
        
        full_content = "\n\n".join(slides_content)
        
        gemini_prompt = f"""Please analyze and structure this PowerPoint content:

{full_content}

Provide a well-structured summary that:
1. Identifies key topics and themes
2. Organizes information logically
3. Highlights important facts, statistics, and findings
4. Makes the content searchable and useful for Q&A
5. Preserves important numerical data and scientific facts

Format as clear, coherent text suitable for document analysis and retrieval."""
        
        enhanced_content = ask_gemini(gemini_prompt, temp=0.3)
        
        if enhanced_content and "Error:" not in enhanced_content:
            logger.info("✓ PowerPoint content enhanced with Gemini")
            return enhanced_content
        else:
            logger.warning("⚠️ Gemini enhancement failed, using raw content")
            return full_content
        
    except Exception as e:
        logger.error(f"Error extracting PowerPoint content: {str(e)}")
        return f"Error extracting PowerPoint content: {str(e)}"

def process_document(filepath, filename):
    ext = filename.rsplit('.', 1)[1].lower()
    
    try:
        if ext in ['pptx', 'ppt']:
            content = extract_ppt_content(filepath)
            if content.startswith("Error"):
                return None, content
            
            from langchain.schema import Document
            
            chunks = []
            sections = content.split('--- Slide')
            for i, section in enumerate(sections):
                if section.strip():
                    if i > 0:
                        section = '--- Slide' + section
                    
                    if len(section) > 2000:
                        sub_chunks = section.split('\n\n')
                        for j, sub_chunk in enumerate(sub_chunks):
                            if sub_chunk.strip():
                                chunks.append(sub_chunk.strip())
                    else:
                        chunks.append(section.strip())
            
            pages = []
            for i, chunk in enumerate(chunks):
                if chunk.strip():
                    doc = Document(
                        page_content=chunk.strip(),
                        metadata={
                            "page": i + 1, 
                            "source": filename, 
                            "type": "presentation", 
                            "file_type": ext,
                            "chunk_id": i
                        }
                    )
                    pages.append(doc)
            
            logger.info(f"✓ PowerPoint processed: {len(pages)} chunks created")
            return pages, None
            
        elif ext in LOADERS:
            loader = LOADERS[ext](filepath, encoding='utf-8') if ext == 'txt' else LOADERS[ext](filepath)
            pages = loader.load()
            if not pages:
                return None, "No content extracted from document"
            
            for i, doc in enumerate(pages):
                doc.metadata = {"page": i + 1, "source": filename, "type": "document", "file_type": ext}
            return pages, None
        else:
            return None, f"Unsupported format: {ext}"
            
    except Exception as e:
        logger.error(f"Error processing {filename}: {str(e)}")
        return None, f"Error processing {filename}: {str(e)}"

def get_document_context(question, k=5):
    try:
        db = Chroma(persist_directory=VECTOR_DB_DIR, embedding_function=HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL))
        docs = db.similarity_search(question, k=k)
        if not docs:
            return None, []
        
        context_parts, sources = [], set()
        for i, doc in enumerate(docs):
            context_parts.append(f"--- Section {i+1} ---")
            context_parts.append(doc.page_content.strip())
            if hasattr(doc, 'metadata') and 'source' in doc.metadata:
                sources.add(doc.metadata['source'])
        
        return "\n\n".join(context_parts), list(sources)
    except Exception as e:
        logger.error(f"Error getting document context: {str(e)}")
        return None, []

def get_numeric_columns(df):
    return df.select_dtypes(include=[np.number]).columns.tolist()

def suggest_chart_columns(df, question):
    numeric_cols = get_numeric_columns(df)
    categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
    
    if not numeric_cols:
        return None, None, None, "No numeric columns found"
    
    # Enhanced chart suggestion logic based on question content
    question_lower = question.lower()
    
    if 'pie' in question_lower and categorical_cols:
        return 'pie', categorical_cols[0], numeric_cols[0], "Pie chart showing distribution"
    elif 'line' in question_lower and len(numeric_cols) >= 2:
        return 'line', numeric_cols[0], numeric_cols[1], "Line chart showing trends"
    elif 'scatter' in question_lower and len(numeric_cols) >= 2:
        return 'scatter', numeric_cols[0], numeric_cols[1], "Scatter plot showing correlation"
    elif 'histogram' in question_lower or 'distribution' in question_lower:
        return 'histogram', None, numeric_cols[0], "Histogram showing distribution"
    elif categorical_cols and numeric_cols:
        return 'bar', categorical_cols[0], numeric_cols[0], "Bar chart with categorical and numeric data"
    elif len(numeric_cols) >= 2:
        return 'scatter', numeric_cols[0], numeric_cols[1], "Scatter plot with two numeric columns"
    else:
        return 'histogram', None, numeric_cols[0], "Distribution histogram"

def generate_smart_chart(df, filename, question):
    """Generate chart with improved error handling and better visualization"""
    chart_type, x_col, y_col, reasoning = suggest_chart_columns(df, question)
    if not chart_type:
        return None, reasoning
    
    try:
        chart_id = uuid.uuid4().hex
        
        # Set up the plot with better styling
        plt.figure(figsize=(12, 8))
        plt.style.use('dark_background')
        
        # Color palette
        colors = ['#667eea', '#764ba2', '#f093fb', '#f5576c', '#4facfe', '#00f2fe']
        
        if chart_type == 'bar' and x_col and y_col:
            if df[x_col].dtype == 'object':
                # Group and aggregate data
                grouped = df.groupby(x_col)[y_col].sum().sort_values(ascending=False).head(10)
                bars = plt.bar(range(len(grouped)), grouped.values, color=colors[0], alpha=0.8)
                plt.xticks(range(len(grouped)), grouped.index, rotation=45, ha='right')
                
                # Add value labels on bars
                for bar in bars:
                    height = bar.get_height()
                    plt.text(bar.get_x() + bar.get_width()/2., height,
                            f'{height:.1f}', ha='center', va='bottom')
            else:
                plt.bar(df[x_col], df[y_col], color=colors[0], alpha=0.8)
            
            plt.xlabel(x_col, fontsize=12)
            plt.ylabel(y_col, fontsize=12)
            plt.title(f'{y_col} by {x_col}', fontsize=14, fontweight='bold')
            
        elif chart_type == 'scatter' and x_col and y_col:
            plt.scatter(df[x_col], df[y_col], alpha=0.7, color=colors[0], s=60)
            plt.xlabel(x_col, fontsize=12)
            plt.ylabel(y_col, fontsize=12)
            plt.title(f'{y_col} vs {x_col}', fontsize=14, fontweight='bold')
            
            # Add trend line
            z = np.polyfit(df[x_col].dropna(), df[y_col].dropna(), 1)
            p = np.poly1d(z)
            plt.plot(df[x_col], p(df[x_col]), "--", color=colors[1], alpha=0.8)
            
        elif chart_type == 'histogram' and y_col:
            plt.hist(df[y_col].dropna(), bins=20, color=colors[0], alpha=0.7, edgecolor='white')
            plt.xlabel(y_col, fontsize=12)
            plt.ylabel('Frequency', fontsize=12)
            plt.title(f'Distribution of {y_col}', fontsize=14, fontweight='bold')
            
        elif chart_type == 'pie' and x_col and y_col:
            if df[x_col].dtype == 'object':
                pie_data = df.groupby(x_col)[y_col].sum().head(8)  # Top 8 categories
                plt.pie(pie_data.values, labels=pie_data.index, autopct='%1.1f%%', 
                       colors=colors, startangle=90)
                plt.title(f'{y_col} Distribution by {x_col}', fontsize=14, fontweight='bold')
                
        elif chart_type == 'line' and x_col and y_col:
            plt.plot(df[x_col], df[y_col], marker='o', linewidth=2, markersize=6, color=colors[0])
            plt.xlabel(x_col, fontsize=12)
            plt.ylabel(y_col, fontsize=12)
            plt.title(f'{y_col} over {x_col}', fontsize=14, fontweight='bold')
            plt.grid(True, alpha=0.3)
        
        # Improve layout and save
        plt.tight_layout()
        plt.subplots_adjust(bottom=0.15)
        
        # Ensure the plots directory exists
        os.makedirs(PLOTS_FOLDER, exist_ok=True)
        
        chart_path = os.path.join(PLOTS_FOLDER, f"{chart_id}.png")
        plt.savefig(chart_path, dpi=300, bbox_inches='tight', facecolor='#1a1f2e', 
                   edgecolor='none', transparent=False)
        plt.close()
        
        # Verify file was created
        if os.path.exists(chart_path):
            logger.info(f"✓ Chart saved successfully: {chart_path}")
            # Return the full URL for the frontend
            chart_url = f"http://localhost:5000/static/plots/{chart_id}.png"
            return chart_url, f"Generated {chart_type} chart: {reasoning}"
        else:
            logger.error("❌ Chart file was not created")
            return None, "Failed to save chart file"
        
    except Exception as e:
        plt.close()  # Ensure plot is closed on error
        logger.error(f"Error generating chart: {str(e)}")
        return None, f"Error generating chart: {str(e)}"

# Routes
@app.route('/')
def index():
    frontend_path = os.path.abspath(os.path.join(os.path.dirname(__file__), '..', 'frontend'))
    return send_from_directory(frontend_path, 'index.html')

@app.route('/health')
def health_check():
    return jsonify({
        'status': 'healthy',
        'gemini_api': bool(GEMINI_API_KEY),
        'has_documents': has_documents(),
        'has_data': has_data(),
        'data_files': list(excel_data_store.keys()),
        'upload_folder': UPLOAD_FOLDER,
        'plots_folder': PLOTS_FOLDER,
        'vector_db_exists': os.path.exists(VECTOR_DB_DIR)
    })

@app.route('/upload', methods=['POST'])
def upload():
    if 'file' not in request.files:
        return jsonify({'error': 'No file part in request'}), 400
    
    file = request.files['file']
    if not file or file.filename == '':
        return jsonify({'error': 'No file selected'}), 400
    
    if not allowed_file(file.filename):
        return jsonify({'error': 'Unsupported file type'}), 400

    filename = secure_filename(file.filename)
    filepath = os.path.join(UPLOAD_FOLDER, filename)
    
    try:
        file.save(filepath)
        logger.info(f"📁 File saved: {filepath}")
    except Exception as e:
        logger.error(f"Failed to save file: {str(e)}")
        return jsonify({'error': f'Failed to save file: {str(e)}'}), 500

    ext = filename.rsplit('.', 1)[1].lower()

    # Handle data files
    if ext in ['xls', 'xlsx', 'csv']:
        try:
            df = pd.read_csv(filepath) if ext == 'csv' else pd.read_excel(filepath)
            excel_data_store[filename] = df
            logger.info(f"📊 Data file processed: {filename} - {len(df)} rows, {len(df.columns)} columns")
            logger.info(f"📊 Columns: {df.columns.tolist()}")
            logger.info(f"📊 Data types: {df.dtypes.to_dict()}")
            
            return jsonify({
                'message': f'Data file uploaded successfully - {len(df)} rows, {len(df.columns)} columns',
                'filename': filename,
                'columns': df.columns.tolist(),
                'shape': df.shape,
                'data_types': df.dtypes.astype(str).to_dict(),
                'sample_data': df.head(3).to_dict()
            })
        except Exception as e:
            logger.error(f"Failed to process data file: {str(e)}")
            return jsonify({'error': f'Failed to process data file: {str(e)}'}), 500
    
    # Handle document files
    else:
        logger.info("📄 Document upload detected - clearing vector database...")
        clear_vector_database()
        
        pages, error = process_document(filepath, filename)
        if error:
            logger.error(f"Document processing error: {error}")
            return jsonify({'error': error}), 500
        
        try:
            text_splitter = RecursiveCharacterTextSplitter(
                chunk_size=1000, 
                chunk_overlap=100,
                separators=["\n\n", "\n", ".", "!", "?", ",", " ", ""]
            )
            texts = text_splitter.split_documents(pages)
            embeddings = HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL)
            
            db = Chroma.from_documents(texts, embeddings, persist_directory=VECTOR_DB_DIR)
            db.persist()
            
            logger.info(f"✓ Document indexed: {len(pages)} pages into {len(texts)} chunks")
            return jsonify({
                'message': f'Document uploaded and processed. {len(pages)} pages into {len(texts)} chunks.',
                'pages': len(pages),
                'chunks': len(texts),
                'file_type': ext,
                'filename': filename
            })
        except Exception as e:
            logger.error(f"Failed to index document: {str(e)}")
            return jsonify({'error': f'Failed to index document: {str(e)}'}), 500

@app.route('/ask', methods=['POST'])
def ask_question():
    data = request.get_json()
    question = data.get('question', '')
    rag_mode = data.get('rag_mode', True)  # New parameter to control RAG vs Web mode
    
    if not question:
        return jsonify({'error': 'No question provided'}), 400
    
    logger.info(f"❓ Question received: {question}")
    logger.info(f"🔄 Mode: {'RAG' if rag_mode else 'Web Search'}")
    
    is_graph_req = is_graph_request(question)
    is_doc_question = is_document_specific_question(question)
    
    try:
        # RAG Mode - Use uploaded files and documents
        if rag_mode:
            logger.info("📚 RAG Mode: Processing with uploaded files...")
            
            # Handle graph requests with priority
            if is_graph_req and has_data():
                logger.info("📊 Processing graph request...")
                filename, df = next(iter(excel_data_store.items()))
                logger.info(f"📊 Using dataset: {filename} with shape {df.shape}")
                
                chart_url, chart_info = generate_smart_chart(df, filename, question)
                
                if chart_url:
                    logger.info(f"✅ Chart generated successfully: {chart_url}")
                    
                    analysis_prompt = f"""I've created a chart based on: "{question}"
Dataset: {filename} ({df.shape[0]} rows × {df.shape[1]} columns)
Chart: {chart_info}
Columns: {', '.join(df.columns.tolist())}
Sample data:
{df.head(3).to_string()}

Provide insightful analysis of this visualization and explain what the chart shows."""
                    
                    return jsonify({
                        'answer': ask_gemini(analysis_prompt, temp=0.5),
                        'chart_url': chart_url,
                        'chart_info': chart_info,
                        'mode': 'rag_graph_generated'
                    })
                else:
                    logger.error(f"❌ Chart generation failed: {chart_info}")
                    return jsonify({
                        'answer': f"I couldn't generate the chart. {chart_info}. Please check your data format and try again.",
                        'mode': 'rag_graph_failed',
                        'error': chart_info
                    })
            
            # Handle document questions
            elif has_documents():
                context, sources = get_document_context(question, k=5)
                
                if context:
                    rag_prompt = f"""Based on the following document content, answer the question comprehensively:

DOCUMENT CONTENT:
{context}

QUESTION: {question}

Provide a detailed, accurate response using the document information:"""

                    answer = ask_gemini(rag_prompt, temp=0.3)
                    return jsonify({
                        'answer': answer,
                        'sources': sources,
                        'mode': 'rag_document',
                        'context_found': True
                    })
                else:
                    return jsonify({
                        'answer': "No relevant information found in uploaded documents for your question.",
                        'mode': 'rag_no_context',
                        'context_found': False
                    })
            
            # No uploaded files in RAG mode
            else:
                return jsonify({
                    'answer': "RAG mode is ON but no files are uploaded. Please upload documents or data files to use RAG functionality, or switch to Web Search mode for general questions.",
                    'mode': 'rag_no_files'
                })
        
        # Web Search Mode - Use Gemini with web context
        else:
            logger.info("🌐 Web Search Mode: Processing with web search...")
            
            # Use Gemini with web search context for current information
            web_enhanced_prompt = f"""Please provide a comprehensive, up-to-date answer to this question: {question}

Use your knowledge and reasoning to provide the most accurate and current information available."""
            
            answer = ask_gemini(web_enhanced_prompt, temp=0.7, include_web_context=True, web_query=question)
            
            return jsonify({
                'answer': answer,
                'mode': 'web_search',
                'web_enhanced': True
            })
        
    except Exception as e:
        logger.error(f"Error processing question: {str(e)}")
        return jsonify({
            'answer': f"Error processing question: {str(e)}",
            'mode': 'error'
        })

@app.route('/static/plots/<filename>')
def serve_plot(filename):
    """Serve generated plot images"""
    try:
        return send_from_directory(PLOTS_FOLDER, filename)
    except Exception as e:
        logger.error(f"Error serving plot {filename}: {str(e)}")
        return jsonify({'error': 'Plot not found'}), 404

@app.route('/clear', methods=['POST'])
def clear_data():
    """Clear all uploaded data and documents"""
    try:
        # Clear vector database
        clear_vector_database()
        
        # Clear data store
        excel_data_store.clear()
        
        # Clear upload folder
        if os.path.exists(UPLOAD_FOLDER):
            for filename in os.listdir(UPLOAD_FOLDER):
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                try:
                    if os.path.isfile(file_path):
                        os.unlink(file_path)
                except Exception as e:
                    logger.error(f"Error deleting {file_path}: {str(e)}")
        
        # Clear plots folder
        if os.path.exists(PLOTS_FOLDER):
            for filename in os.listdir(PLOTS_FOLDER):
                file_path = os.path.join(PLOTS_FOLDER, filename)
                try:
                    if os.path.isfile(file_path):
                        os.unlink(file_path)
                except Exception as e:
                    logger.error(f"Error deleting {file_path}: {str(e)}")
        
        logger.info("🧹 All data cleared successfully")
        return jsonify({'message': 'All data cleared successfully'})
        
    except Exception as e:
        logger.error(f"Error clearing data: {str(e)}")
        return jsonify({'error': f'Error clearing data: {str(e)}'}), 500

@app.route('/toggle-mode', methods=['POST'])
def toggle_mode():
    """Endpoint to get current mode status"""
    data = request.get_json()
    rag_mode = data.get('rag_mode', True)
    
    # Return current status
    return jsonify({
        'rag_mode': rag_mode,
        'has_documents': has_documents(),
        'has_data': has_data(),
        'message': f"Mode set to {'RAG (Files)' if rag_mode else 'Web Search'}"
    })

@app.route('/get-data-info', methods=['GET'])
def get_data_info():
    """Get information about uploaded data files"""
    if not excel_data_store:
        return jsonify({
            'has_data': False,
            'message': 'No data files uploaded'
        })
    
    data_info = {}
    for filename, df in excel_data_store.items():
        data_info[filename] = {
            'shape': df.shape,
            'columns': df.columns.tolist(),
            'numeric_columns': get_numeric_columns(df),
            'data_types': df.dtypes.astype(str).to_dict(),
            'sample': df.head(2).to_dict()
        }
    
    return jsonify({
        'has_data': True,
        'files': data_info,
        'total_files': len(excel_data_store)
    })

@app.route('/get-document-info', methods=['GET'])
def get_document_info():
    """Get information about uploaded documents"""
    try:
        if not has_documents():
            return jsonify({
                'has_documents': False,
                'message': 'No documents uploaded'
            })
        
        db = Chroma(persist_directory=VECTOR_DB_DIR, 
                   embedding_function=HuggingFaceEmbeddings(model_name=EMBEDDING_MODEL))
        
        # Get a sample to understand document structure
        sample_docs = db.similarity_search("sample", k=3)
        
        doc_info = {
            'has_documents': True,
            'total_chunks': len(sample_docs) if sample_docs else 0,
            'sample_content': []
        }
        
        for i, doc in enumerate(sample_docs[:2]):
            doc_info['sample_content'].append({
                'chunk_id': i,
                'content_preview': doc.page_content[:200] + "..." if len(doc.page_content) > 200 else doc.page_content,
                'metadata': doc.metadata
            })
        
        return jsonify(doc_info)
        
    except Exception as e:
        logger.error(f"Error getting document info: {str(e)}")
        return jsonify({
            'has_documents': False,
            'error': str(e)
        })

@app.route('/analyze-data', methods=['POST'])
def analyze_data():
    """Provide statistical analysis of uploaded data"""
    if not excel_data_store:
        return jsonify({'error': 'No data files uploaded'}), 400
    
    data = request.get_json()
    filename = data.get('filename')
    
    try:
        if filename and filename in excel_data_store:
            df = excel_data_store[filename]
        else:
            # Use first available dataset
            filename, df = next(iter(excel_data_store.items()))
        
        analysis = {
            'filename': filename,
            'shape': df.shape,
            'columns': df.columns.tolist(),
            'numeric_columns': get_numeric_columns(df),
            'missing_values': df.isnull().sum().to_dict(),
            'data_types': df.dtypes.astype(str).to_dict()
        }
        
        # Basic statistics for numeric columns
        numeric_cols = get_numeric_columns(df)
        if numeric_cols:
            analysis['statistics'] = df[numeric_cols].describe().to_dict()
        
        # Generate analysis with Gemini
        analysis_prompt = f"""Analyze this dataset:

Dataset: {filename}
Shape: {df.shape[0]} rows × {df.shape[1]} columns
Columns: {', '.join(df.columns.tolist())}
Numeric columns: {numeric_cols}
Missing values: {dict(df.isnull().sum())}

Sample data:
{df.head(5).to_string()}

Provide insights about:
1. Data quality and completeness
2. Key patterns or trends visible
3. Potential analysis opportunities
4. Recommended visualizations
5. Data preparation suggestions"""

        ai_analysis = ask_gemini(analysis_prompt, temp=0.5)
        analysis['ai_insights'] = ai_analysis
        
        return jsonify(analysis)
        
    except Exception as e:
        logger.error(f"Error analyzing data: {str(e)}")
        return jsonify({'error': f'Error analyzing data: {str(e)}'}), 500

@app.route('/suggest-questions', methods=['POST'])
def suggest_questions():
    """Suggest relevant questions based on uploaded content"""
    suggestions = []
    
    try:
        # Suggestions for data files
        if excel_data_store:
            for filename, df in excel_data_store.items():
                numeric_cols = get_numeric_columns(df)
                categorical_cols = df.select_dtypes(include=['object', 'category']).columns.tolist()
                
                if numeric_cols and categorical_cols:
                    suggestions.extend([
                        f"Create a bar chart of {numeric_cols[0]} by {categorical_cols[0]}",
                        f"Show me the distribution of {numeric_cols[0]}",
                        f"Create a scatter plot of {numeric_cols[0]} vs {numeric_cols[1] if len(numeric_cols) > 1 else numeric_cols[0]}",
                        f"Analyze the relationship between {categorical_cols[0]} and {numeric_cols[0]}",
                        f"What are the key insights from the {filename} data?"
                    ])
                elif numeric_cols:
                    suggestions.extend([
                        f"Show me the distribution of {numeric_cols[0]}",
                        f"Create a histogram of {numeric_cols[0]}",
                        f"What are the statistics for {numeric_cols[0]}?"
                    ])
        
        # Suggestions for documents
        if has_documents():
            suggestions.extend([
                "Summarize the main points from the document",
                "What are the key findings mentioned?",
                "Explain the methodology described in the document",
                "What recommendations are provided?",
                "Analyze the conclusions presented"
            ])
        
        # General suggestions if no files
        if not excel_data_store and not has_documents():
            suggestions.extend([
                "What is machine learning?",
                "Explain artificial intelligence concepts",
                "How does data analysis work?",
                "What are the latest trends in technology?",
                "Tell me about data visualization best practices"
            ])
        
        return jsonify({
            'suggestions': suggestions[:8],  # Limit to 8 suggestions
            'has_data': len(excel_data_store) > 0,
            'has_documents': has_documents()
        })
        
    except Exception as e:
        logger.error(f"Error generating suggestions: {str(e)}")
        return jsonify({
            'suggestions': ["Ask me anything!"],
            'error': str(e)
        })

# Error handlers
@app.errorhandler(404)
def not_found_error(error):
    return jsonify({'error': 'Endpoint not found'}), 404

@app.errorhandler(500)
def internal_error(error):
    logger.error(f"Internal server error: {str(error)}")
    return jsonify({'error': 'Internal server error'}), 500

@app.errorhandler(413)
def too_large(error):
    return jsonify({'error': 'File too large. Maximum size is 50MB.'}), 413

# Cleanup function
def cleanup_old_files():
    """Clean up old uploaded files and plots"""
    try:
        import time
        current_time = time.time()
        max_age = 24 * 60 * 60  # 24 hours
        
        # Clean old uploads
        if os.path.exists(UPLOAD_FOLDER):
            for filename in os.listdir(UPLOAD_FOLDER):
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                if os.path.isfile(file_path):
                    if current_time - os.path.getctime(file_path) > max_age:
                        try:
                            os.unlink(file_path)
                            logger.info(f"🧹 Cleaned old upload: {filename}")
                        except Exception as e:
                            logger.error(f"Error cleaning {file_path}: {str(e)}")
        
        # Clean old plots
        if os.path.exists(PLOTS_FOLDER):
            for filename in os.listdir(PLOTS_FOLDER):
                file_path = os.path.join(PLOTS_FOLDER, filename)
                if os.path.isfile(file_path):
                    if current_time - os.path.getctime(file_path) > max_age:
                        try:
                            os.unlink(file_path)
                            logger.info(f"🧹 Cleaned old plot: {filename}")
                        except Exception as e:
                            logger.error(f"Error cleaning {file_path}: {str(e)}")
                            
    except Exception as e:
        logger.error(f"Error in cleanup: {str(e)}")

if __name__ == '__main__':
    logger.info("🚀 Starting Flask application...")
    logger.info(f"📁 Upload folder: {UPLOAD_FOLDER}")
    logger.info(f"📊 Plots folder: {PLOTS_FOLDER}")
    logger.info(f"💾 Vector DB folder: {VECTOR_DB_DIR}")
    logger.info(f"🤖 Gemini API: {'✓ Configured' if GEMINI_API_KEY else '❌ Not configured'}")
    
    # Run cleanup on startup
    cleanup_old_files()
    
    # Configure Flask app
    app.config['MAX_CONTENT_LENGTH'] = 50 * 1024 * 1024  # 50MB max file size
    app.config['UPLOAD_FOLDER'] = UPLOAD_FOLDER
    
    # Run the application
    app.run(
        host='0.0.0.0',
        port=5000,
        debug=True,
        threaded=True
    )