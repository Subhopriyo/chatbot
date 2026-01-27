"""
PDF processing service with highlighting support
"""
import os
import logging
from typing import Tuple, List, Optional, Dict
from langchain.schema import Document
from pptx import Presentation

from services.groq_service import grok

logger = logging.getLogger(__name__)

# Document loaders
LOADERS = {}
try:
    from langchain_community.document_loaders import PyPDFLoader, TextLoader, Docx2txtLoader
    LOADERS.update({
        'pdf': PyPDFLoader,
        'txt': TextLoader,
        'docx': Docx2txtLoader
    })
except ImportError:
    logger.warning("Some document loaders unavailable")

def extract_ppt_content(filepath) -> str:
    """Extract and enhance PowerPoint content"""
    try:
        # Convert Path to string if needed
        filepath = str(filepath)
        
        if not os.path.exists(filepath):
            return f"Error: File not found at {filepath}"
        
        file_size = os.path.getsize(filepath)
        if file_size > 50 * 1024 * 1024:
            return "Error: File too large (>50MB)"
        
        logger.info(f"📄 Processing PowerPoint: {filepath}")
        
        try:
            prs = Presentation(filepath)
        except Exception as e:
            return f"Error: Unable to open PowerPoint - {str(e)}"
        
        slides_content = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = [f"--- Slide {i+1} ---"]
            text_found = False
            
            for shape in slide.shapes:
                try:
                    # Extract text
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                        text_found = True
                    
                    # Extract tables
                    if hasattr(shape, "table"):
                        table_text = []
                        for row in shape.table.rows:
                            row_text = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                            if row_text:
                                table_text.append(" | ".join(row_text))
                        if table_text:
                            slide_text.extend(table_text)
                            text_found = True
                            
                except Exception as e:
                    logger.warning(f"Shape processing error on slide {i+1}: {str(e)}")
                    continue
            
            if not text_found:
                slide_text.append("[No text content found in this slide]")
            
            slides_content.append("\n".join(slide_text))
        
        if not slides_content:
            return "Error: No slides found"
        
        full_content = "\n\n".join(slides_content)
        
        # Enhance with Groq
        enhancement_prompt = f"""Analyze and structure this PowerPoint content:

{full_content}

Provide a well-structured summary that:
1. Identifies key topics and themes
2. Organizes information logically
3. Highlights important facts and statistics
4. Makes content searchable for Q&A
5. Preserves numerical data

Format as clear, coherent text for document analysis."""

        enhanced = grok.simple_prompt(enhancement_prompt, temperature=0.3)
        
        if enhanced and not enhanced.startswith("⚠️"):
            logger.info("✓ PowerPoint enhanced with Groq")
            return enhanced
        else:
            logger.warning("⚠️ Using raw PowerPoint content")
            return full_content
            
    except Exception as e:
        logger.error(f"PowerPoint extraction error: {str(e)}")
        return f"Error: {str(e)}"

def process_document(filepath, filename: str) -> Tuple[Optional[List[Document]], Optional[str]]:
    """
    Process document and return pages
    
    Args:
        filepath: File path (can be Path object or string)
        filename: Original filename
    
    Returns:
        (pages, error_message)
    """
    # Convert Path to string if needed
    filepath = str(filepath)
    
    ext = filename.rsplit('.', 1)[1].lower()
    
    try:
        # PowerPoint files
        if ext in ['pptx', 'ppt']:
            content = extract_ppt_content(filepath)
            
            if content.startswith("Error"):
                return None, content
            
            # Split into chunks
            chunks = []
            sections = content.split('--- Slide')
            
            for i, section in enumerate(sections):
                if section.strip():
                    if i > 0:
                        section = '--- Slide' + section
                    
                    # Further split long sections
                    if len(section) > 2000:
                        sub_chunks = section.split('\n\n')
                        chunks.extend([c for c in sub_chunks if c.strip()])
                    else:
                        chunks.append(section.strip())
            
            # Create documents
            pages = []
            for i, chunk in enumerate(chunks):
                if chunk.strip():
                    doc = Document(
                        page_content=chunk.strip(),
                        metadata={
                            "page": i + 1,
                            "source": filename,
                            "type": "presentation",
                            "file_type": ext,
                            "chunk_id": i
                        }
                    )
                    pages.append(doc)
            
            logger.info(f"✓ PowerPoint processed: {len(pages)} chunks")
            return pages, None
        
        # Other documents
        elif ext in LOADERS:
            loader = LOADERS[ext](filepath)
            pages = loader.load()
            
            if not pages:
                return None, "No content extracted"
            
            # Add metadata
            for i, doc in enumerate(pages):
                doc.metadata = {
                    "page": i + 1,
                    "source": filename,
                    "type": "document",
                    "file_type": ext
                }
            
            return pages, None
        
        else:
            return None, f"Unsupported format: {ext}"
            
    except Exception as e:
        logger.error(f"Document processing error: {str(e)}")
        return None, f"Error: {str(e)}"

def extract_pdf_page_text(filepath, page_number: int) -> Optional[str]:
    """Extract text from specific PDF page"""
    try:
        import fitz  # PyMuPDF
        
        # Convert Path to string if needed
        filepath = str(filepath)
        
        doc = fitz.open(filepath)
        if page_number < 1 or page_number > len(doc):
            return None
        
        page = doc[page_number - 1]
        text = page.get_text()
        doc.close()
        
        return text
        
    except Exception as e:
        logger.error(f"PDF text extraction error: {str(e)}")
        return None

def find_text_in_pdf(filepath, search_text: str) -> List[Dict]:
    """
    Find text locations in PDF for highlighting
    
    Args:
        filepath: PDF file path (can be Path object or string)
        search_text: Text to search for
    
    Returns:
        List of {page, rect} dicts
    """
    try:
        import fitz
        
        # Convert Path to string if needed
        filepath = str(filepath)
        
        doc = fitz.open(filepath)
        results = []
        
        for page_num in range(len(doc)):
            page = doc[page_num]
            text_instances = page.search_for(search_text)
            
            for rect in text_instances:
                results.append({
                    'page': page_num + 1,
                    'rect': {
                        'x0': rect.x0,
                        'y0': rect.y0,
                        'x1': rect.x1,
                        'y1': rect.y1
                    }
                })
        
        doc.close()
        return results
        
    except Exception as e:
        logger.error(f"PDF search error: {str(e)}")
        return []